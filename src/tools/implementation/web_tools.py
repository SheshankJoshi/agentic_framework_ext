from langchain_core.tools import Tool
from langchain_google_community import GoogleSearchAPIWrapper
from functools import partial
import os
from pptx import Presentation

# # -- should be removed later
# from dotenv import load_dotenv
# load_dotenv()
# # -- should be removed later

def google_search_web(num_results=5):
    search = GoogleSearchAPIWrapper()
    func = partial(search.results, num_results=num_results)
    return func

# Function: Generate a PowerPoint presentation with a references slide
def generate_presentation(title: str, content: str, references: list):
    """"Generates a PowerPoint presentation with provided title, content and references."""
    prs = Presentation()

    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = content

    # Create a references slide
    ref_slide_layout = prs.slide_layouts[1] if len(
        prs.slide_layouts) > 1 else title_slide_layout
    ref_slide = prs.slides.add_slide(ref_slide_layout)
    ref_slide.shapes.title.text = "References"
    ref_text = "\n".join(references) if references else "No references found."
    if len(ref_slide.placeholders) > 1:
        ref_slide.placeholders[1].text = ref_text

    # Save the presentation
    ppt_path = os.path.join(os.getcwd(), "presentation.pptx")
    prs.save(ppt_path)
    return f"Presentation saved to: {ppt_path}"

# Function: Process references (for example, formatting the list)


def process_references(references: list):
    if not references:
        return "No references to process."
    ref_output = "Processed References:\n" + "\n".join(references)
    return ref_output


# Wrap our functions as tools for the agent.
# Note: We use the robust web_search tool imported from basic_tools.

